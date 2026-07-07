"""把 raw 下的源文档（docx/doc/pdf/pptx）转成 markdown。

输入：knowledge-base/raw/<file> 下的源文档
输出：knowledge-base/raw/markdown/<同名>.md

支持格式：
- .docx → python-docx
- .doc  → docx2txt
- .pdf  → pdfplumber（按页分割）
- .pptx → python-pptx（slide + notes）

不做 OCR、不还原复杂版式；目标是让下游 LLM/正则能读到结构化文本。
已存在输出默认跳过，--force 覆盖。
"""
from __future__ import annotations

import argparse
import io
import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
DEFAULT_OUTPUT_DIR = PROJECT_ROOT / "knowledge-base" / "raw" / "markdown"


def _stdout_utf8() -> None:
    """Windows 终端默认 cp936，强制 utf-8 输出。"""
    if hasattr(sys.stdout, "buffer"):
        sys.stdout = io.TextIOWrapper(
            sys.stdout.buffer, encoding="utf-8", line_buffering=True
        )


def convert_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append("|".join(cells))
    return "\n\n".join(parts)


def convert_doc(path: Path) -> str:
    import docx2txt

    return docx2txt.process(str(path)) or ""


def convert_pdf(path: Path) -> str:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            parts.append(f"--- Page {i} ---")
            text = page.extract_text() or ""
            if text.strip():
                parts.append(text)
    return "\n\n".join(parts)


def convert_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        seen: set[str] = set()
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs).strip()
                if t and t not in seen:
                    seen.add(t)
                    parts.append(t)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Notes] {notes}")
    return "\n\n".join(parts)


CONVERTERS = {
    ".docx": convert_docx,
    ".doc": convert_doc,
    ".pdf": convert_pdf,
    ".pptx": convert_pptx,
}


def main() -> int:
    _stdout_utf8()
    p = argparse.ArgumentParser(description="把源文档转成 markdown")
    p.add_argument("input", type=Path, help="knowledge-base/raw/ 下的源文件")
    p.add_argument("--out-dir", type=Path, default=DEFAULT_OUTPUT_DIR)
    p.add_argument("--force", action="store_true", help="覆盖已存在的输出")
    args = p.parse_args()

    if not args.input.exists():
        print(f"找不到输入: {args.input}")
        return 1
    suffix = args.input.suffix.lower()
    if suffix not in CONVERTERS:
        print(f"不支持的格式: {suffix}（支持 .docx/.doc/.pdf/.pptx）")
        return 1

    out_path = args.out_dir / (args.input.stem + ".md")
    if out_path.exists() and not args.force:
        print(f"已存在，跳过: {out_path}\n（用 --force 覆盖）")
        return 0

    print(f"转换 {args.input.name} ...")
    try:
        text = CONVERTERS[suffix](args.input)
    except Exception as e:
        print(f"转换失败: {e}")
        return 2
    if not text.strip():
        print(f"⚠️ 输出为空: {args.input.name}")
        return 2

    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(text, encoding="utf-8")
    print(f"✓ 写入 {out_path}（{len(text)} 字）")
    return 0


if __name__ == "__main__":
    sys.exit(main())
