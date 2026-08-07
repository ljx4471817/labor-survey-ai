"""Extract per-page text + real images from a PPTX into a structured dir.

Usage:
    python extract_pptx.py <src.pptx> <out_dir>

Produces:
    <out_dir>/pages.json   # [{page,title,body,notes,images:[relpaths]}]
    <out_dir>/pages.md     # human-readable per-page dump
    <out_dir>/images/page_NN/*.jpeg|*.png|...

CRITICAL (Windows): zip-internal paths use forward slashes. Never use
os.path.normpath on them (it back-slashes and silently drops images);
use posixpath so `in z.namelist()` matches.
"""
import zipfile, re, os, json, posixpath, sys
from pptx import Presentation

SRC = sys.argv[1] if len(sys.argv) > 1 else r"源文件.pptx"
OUT = sys.argv[2] if len(sys.argv) > 2 else r"pptx_extract"
IMG = os.path.join(OUT, "images")
os.makedirs(IMG, exist_ok=True)

z = zipfile.ZipFile(SRC)
prs = Presentation(SRC)


def media_for_slide(n):
    sxml = z.read(f"ppt/slides/slide{n}.xml").decode('utf-8', 'ignore')
    embeds = re.findall(r'r:embed="([^"]+)"', sxml)  # image rel ids
    rel = f"ppt/slides/_rels/slide{n}.xml.rels"
    try:
        rtxt = z.read(rel).decode('utf-8', 'ignore')
    except KeyError:
        return []
    res = []
    for e in embeds:
        m = re.search(rf'Id="{re.escape(e)}"[^>]*Target="([^"]+)"', rtxt)
        if m:
            t = m.group(1)
            full = posixpath.normpath(posixpath.join("ppt/slides", t))  # MUST be posixpath
            res.append(full)
    return res


pages = []
for slide in prs.slides:
    n = int(re.findall(r'\d+', slide.part.partname)[-1])
    texts = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t:
                texts.append(t)
    title = texts[0] if texts else ""
    body = "\n".join(texts[1:]) if len(texts) > 1 else ""
    notes = ""
    try:
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        pass
    media = media_for_slide(n)
    img_files = []
    pdir = os.path.join(IMG, f"page_{n:02d}")
    if media:
        os.makedirs(pdir, exist_ok=True)
    for m in media:
        if m in z.namelist():
            data = z.read(m)
            fname = os.path.basename(m)
            dest = os.path.join(pdir, fname)
            with open(dest, "wb") as f:
                f.write(data)
            img_files.append(posixpath.relpath(dest, OUT).replace("\\", "/"))
        else:
            print("WARN missing in zip:", m)
    pages.append({"page": n, "title": title, "body": body, "notes": notes, "images": img_files})

with open(os.path.join(OUT, "pages.json"), "w", encoding="utf-8") as f:
    json.dump(pages, f, ensure_ascii=False, indent=2)

with open(os.path.join(OUT, "pages.md"), "w", encoding="utf-8") as f:
    for p in pages:
        f.write(f"## 第 {p['page']} 页\n\n")
        if p["title"]:
            f.write(f"**标题**: {p['title']}\n\n")
        if p["body"]:
            f.write(p["body"] + "\n\n")
        if p["notes"]:
            f.write(f"**备注**: {p['notes']}\n\n")
        for img in p["images"]:
            f.write(f"![{img}]({img})\n\n")

total_imgs = sum(len(p["images"]) for p in pages)
print("pages:", len(pages), "| images extracted:", total_imgs)
