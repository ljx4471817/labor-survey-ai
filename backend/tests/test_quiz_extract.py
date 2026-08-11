# -*- coding: utf-8 -*-
"""quiz_extract 纯函数与编排测试（PRD v3 10.1）。"""
from __future__ import annotations

import pytest

from app.services import quiz_extract as qe


def test_segment_notice_groups_sections():
    notice = "8月工作提示\n一、审核要点\n将家务劳动者误判为就业。\n（二）问卷要点\n调查参考周为8月3-9日。"
    segs = qe.segment_notice(notice)
    sections = [s["section"] for s in segs]
    assert sections == ["其它", "审核要点", "问卷要点"]
    assert "将家务劳动者误判为就业" in segs[1]["text"]


def test_segment_notice_title_not_header():
    """文档标题含「工作提示」不得被误判为章节标题。"""
    segs = qe.segment_notice("8月劳动力调查工作提示\n一、审核要点\n要点内容")
    assert segs[0]["section"] == "其它"
    assert segs[1]["section"] == "审核要点"


def test_strip_json_fence():
    raw = "```json\n{\"a\": 1}\n```"
    assert qe.strip_json_fence(raw) == '{"a": 1}'
    assert qe.strip_json_fence('{"a": 1}') == '{"a": 1}'


def test_parse_llm_json_valid_and_markdown():
    raw = '```json\n[{"section": "审核要点", "content": "c", "common_error": "e", "suggest_quiz": true}]\n```'
    data = qe.parse_llm_json(raw)
    assert data == [{"section": "审核要点", "content": "c", "common_error": "e", "suggest_quiz": True}]


def test_parse_llm_json_invalid_returns_none():
    assert qe.parse_llm_json("不是 JSON") is None
    assert qe.parse_llm_json("") is None


def test_parse_keypoints_returns_none_on_invalid_json():
    assert qe.parse_keypoints("不是 JSON") is None
    assert qe.parse_keypoints("") is None


def test_normalize_keypoints_skips_invalid():
    items = [
        {"section": "审核要点", "content": "要点A", "suggest_quiz": True},
        {"section": "x", "content": "  "},  # 空 content，跳过
        "not dict",  # 非 dict，跳过
    ]
    out2 = qe.normalize_keypoints(items)
    assert len(out2) == 1
    assert out2[0]["content"] == "要点A"


def test_find_source_quote_matches_sentence():
    text = "将家务劳动者误判为就业是本月常见错误。时间安排见附件。"
    quote = qe.find_source_quote("将家务劳动者误判为就业", text)
    assert quote == "将家务劳动者误判为就业是本月常见错误"


def test_run_extraction_with_fake_llm():
    def fake_llm(messages):
        return '[{"section": "审核要点", "content": "家务劳动者无收入应判为非劳动力", "common_error": "误判为就业", "suggest_quiz": true}]'
    notice = "一、审核要点\n将家务劳动者误判为就业是本月常见错误。"
    out = qe.run_extraction(notice, fake_llm)
    assert len(out) == 1
    assert out[0]["section"] == "审核要点"
    assert out[0]["source_quote"]  # 应匹配到原文句子


def test_run_extraction_retries_on_invalid_json():
    calls = {"n": 0}

    def fake_llm(messages):
        calls["n"] += 1
        if calls["n"] == 1:
            return "抱歉，我无法输出 JSON"
        return '[{"section": "审核要点", "content": "要点", "suggest_quiz": true}]'

    out = qe.run_extraction("一、审核要点\n要点内容", fake_llm)
    assert len(out) == 1
    assert calls["n"] >= 2


def test_run_extraction_fails_after_retries():
    def fake_llm(messages):
        return "still not json"

    with pytest.raises(ValueError):
        qe.run_extraction("一、审核要点\n要点内容", fake_llm)


def test_segment_notice_cn_numbered_headers():
    """真实通知结构：一、/（一）中文序号标题应被切分，1. 阿拉伯序号内容行不切分。"""
    notice = (
        "关于做好2026年8月份劳动力调查工作的通知\n"
        "一、调查时间安排\n"
        "调查标准时点：8月10日零时。\n"
        "（一）进一步做好数据审核监测\n"
        "1.加强数据指标核查。①表尾备注与指标填报逻辑不一致性。\n"
        "2.加强汇总指标监测。\n"
        "（二）其他事项\n"
        "注意调查安全。\n"
    )
    segs = qe.segment_notice(notice)
    sections = [s["section"] for s in segs]
    assert sections == ["其它", "调查时间安排", "进一步做好数据审核监测", "其他事项"]
    # 阿拉伯序号行并入所属章节，不单独成段
    assert "1.加强数据指标核查" in segs[2]["text"]


def test_extract_doc_text_dispatches_by_suffix(tmp_path, monkeypatch):
    import app.services.quiz_extract as qe

    calls = {}
    monkeypatch.setattr(qe, "extract_docx_text", lambda p: (calls.__setitem__("docx", calls.get("docx", 0) + 1) or "DOCX"))
    monkeypatch.setattr(qe, "_extract_legacy_doc_text", lambda p: (calls.__setitem__("legacy", calls.get("legacy", 0) + 1) or "LEGACY"))
    for name in ("a.docx", "b.doc", "c.wps", "d.DOC"):
        p = tmp_path / name
        p.write_text("x")
        qe.extract_doc_text(str(p))
    assert calls == {"docx": 1, "legacy": 3}


def test_legacy_doc_com_success(tmp_path, monkeypatch):
    import win32com.client

    import app.services.quiz_extract as qe

    p = tmp_path / "a.doc"
    p.write_text("x")

    class FakeDoc:
        Content = type("C", (), {"Text": "这是测试内容\n第二行"})()

        def Close(self, save):
            pass

    class FakeDocs:
        def Open(self, path, ReadOnly=False):
            return FakeDoc()

    class FakeApp:
        def __init__(self):
            self.Visible = True
            self.Documents = FakeDocs()

    calls = []

    def fake_dispatch(prog_id):
        calls.append(prog_id)
        return FakeApp()

    monkeypatch.setattr(win32com.client, "Dispatch", fake_dispatch)
    text = qe._extract_legacy_doc_text(str(p))
    assert "这是测试内容" in text
    assert calls == ["Word.Application"]


def test_legacy_doc_com_failure_raises_clear_error(tmp_path, monkeypatch):
    import win32com.client

    import app.services.quiz_extract as qe

    p = tmp_path / "a.wps"
    p.write_text("x")

    def fake_dispatch(prog_id):
        raise Exception("cannot start " + prog_id)

    monkeypatch.setattr(win32com.client, "Dispatch", fake_dispatch)
    with pytest.raises(RuntimeError, match="无法解析"):
        qe._extract_legacy_doc_text(str(p))


def test_try_soffice_missing_returns_none(tmp_path, monkeypatch):
    import shutil

    import app.services.quiz_extract as qe

    monkeypatch.setattr(shutil, "which", lambda *a: None)
    p = tmp_path / "a.doc"
    p.write_text("x")
    assert qe._try_soffice(str(p)) is None


def test_try_soffice_success(tmp_path, monkeypatch):
    import shutil
    import subprocess
    import tempfile
    import types

    import app.services.quiz_extract as qe

    p = tmp_path / "a.doc"
    p.write_text("x")
    out = tmp_path / "out"
    out.mkdir()
    (out / "a.txt").write_text("这是转换后的内容", encoding="utf-8")
    monkeypatch.setattr(shutil, "which", lambda *a: "soffice")
    calls = {}

    def fake_run(cmd, capture_output=True, timeout=120):
        calls["cmd"] = cmd
        r = types.SimpleNamespace(returncode=0)
        return r

    monkeypatch.setattr(subprocess, "run", fake_run)
    monkeypatch.setattr(tempfile, "TemporaryDirectory", lambda: _FakeTmp(str(out)))
    text = qe._try_soffice(str(p))
    assert text == "这是转换后的内容"
    assert "--convert-to" in calls["cmd"]


class _FakeTmp:
    def __init__(self, d):
        self.d = d

    def __enter__(self):
        return self.d

    def __exit__(self, *a):
        pass
# --- 多格式文本提取（PRD v5：word/ppt/pdf）---


def _make_pdf(path: str, text: str = "Hello Survey PDF") -> None:
    """构造最小合法 PDF（Helvetica 文本层），供 pdfplumber 提取。"""
    from pathlib import Path
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
    ]
    stream = f"BT /F1 18 Tf 100 700 Td ({text}) Tj ET".encode()
    objs.append(b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream")
    objs.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")
    out = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for i, o in enumerate(objs, 1):
        offsets.append(len(out))
        out += f"{i} 0 obj\n".encode() + o + b"\nendobj\n"
    xref_pos = len(out)
    out += f"xref\n0 {len(objs) + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for off in offsets[1:]:
        out += f"{off:010d} 00000 n \n".encode()
    out += f"trailer\n<< /Size {len(objs) + 1} /Root 1 0 R >>\nstartxref\n{xref_pos}\n%%EOF".encode()
    Path(path).write_bytes(bytes(out))


def test_extract_pdf_text(tmp_path):
    from app.services.quiz_extract import extract_pdf_text

    p = tmp_path / "材料.pdf"
    _make_pdf(str(p), "Hello Survey PDF")
    text = extract_pdf_text(str(p))
    assert "Hello Survey PDF" in text


def test_extract_pptx_text(tmp_path):
    from app.services.quiz_extract import extract_pptx_text
    from pptx import Presentation
    from pptx.util import Inches

    p = tmp_path / "培训.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    tb.text_frame.text = "新员工上岗培训要点：参考周为8月3-9日"
    slide.notes_slide.notes_text_frame.text = "备注：强调参考周"
    prs.save(p)
    text = extract_pptx_text(str(p))
    assert "参考周" in text and "备注" in text


def test_extract_file_text_dispatches(tmp_path):
    from app.services.quiz_extract import extract_file_text

    p = tmp_path / "材料.pdf"
    _make_pdf(str(p), "Hello Survey PDF")
    assert "Hello Survey PDF" in extract_file_text(str(p))


# --- 提取目标要点数（PRD v5）---


def test_run_extraction_keypoint_count(monkeypatch):
    from app.services import quiz_extract as qe

    captured = {}

    def fake_llm_json(messages, llm_chat_fn, parse_fn, what):
        captured["prompt"] = messages[1]["content"]
        return [{"section": "审核要点", "content": "要点A"}]

    monkeypatch.setattr(qe, "_llm_json", fake_llm_json)
    # 默认 10
    qe.run_extraction("文本", lambda m: "")
    assert "提取 10 个可出题的要点" in captured["prompt"]
    # 指定 3
    qe.run_extraction("文本", lambda m: "", keypoint_count=3)
    assert "提取 3 个可出题的要点" in captured["prompt"]
    # 超上限截断到 30
    qe.run_extraction("文本", lambda m: "", keypoint_count=99)
    assert "提取 30 个可出题的要点" in captured["prompt"]
    # 低于下限抬到 1
    qe.run_extraction("文本", lambda m: "", keypoint_count=0)
    assert "提取 1 个可出题的要点" in captured["prompt"]
