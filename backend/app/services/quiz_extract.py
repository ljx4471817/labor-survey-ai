# -*- coding: utf-8 -*-
"""测验：要点提取服务（docx/pdf/pptx → 文本 → 章节 → LLM 要点 → JSON 解析）。

分层（PRD v3 6.1 / 6.5）：
- 纯函数（本模块，可单测）：extract_docx_text / segment_notice / strip_json_fence /
  parse_llm_json / normalize_keypoints / find_source_quote / parse_keypoints
- IO 编排（依赖注入 llm_chat_fn，测试时 mock）：run_extraction
"""
from __future__ import annotations

import json
import re
from difflib import SequenceMatcher
from pathlib import Path

from app.core.constants import QUIZ_DEFAULT_KEYPOINTS, QUIZ_MAX_KEYPOINTS, QUIZ_RETRY_TIMES

# 章节识别关键词（命中且行短 → 视为新章节标题）
# 允许上传的文档类型：docx 走 python-docx；doc/wps 走本机 Word/WPS COM；
# pdf 走 pdfplumber（文字层）；pptx 走 python-pptx（文本框+备注，图片内容不提取）
ALLOWED_DOC_EXTENSIONS: tuple[str, ...] = (".docx", ".doc", ".wps", ".pdf", ".pptx")


SECTION_HEADER_KEYWORDS: tuple[str, ...] = (
    "审核要点", "问卷要点", "填报口径", "注意事项",
)

KNOWN_SECTIONS: dict[str, str] = {
    "审核要点": "审核要点",
    "问卷要点": "问卷要点",
    "填报口径": "填报口径微调",
    "填报口径微调": "填报口径微调",
}


def _strip_numbering(s: str) -> str:
    """去掉章节号前缀：一、 / 1. / （一）。"""
    return re.sub(r"^[（(]?[一二三四五六七八九十0-9]+[）)）、.\s]+", "", s.strip().strip("：: ."))


def _is_section_header(line: str) -> bool:
    """标题识别：短行且为中文序号开头（一、 / （一）），或恰为已知章节名。

    覆盖真实通知结构（一、调查时间安排 / （一）进一步做好数据审核监测…），
    同时避免把文档标题「…工作提示」或阿拉伯序号内容行（1. xxx）误判为章节。
    """
    s = line.strip().strip("：: .")
    if len(s) > 30:
        return False
    if s in KNOWN_SECTIONS:
        return True
    return bool(re.match(r"^[（(]?[一二三四五六七八九十百]+[）)]?[、.．]?\s*", s))

PROMPT1_SYSTEM = (
    "你是劳动力调查专家。请从给定的文件内容（通知、培训材料、制度文件等）中"
    "提取需要调查员记住并应用的要点，只输出 JSON。"
)

PROMPT1_USER = """从以下文件内容中提取 {keypoint_count} 个可出题的要点。

## 输入
{notice_text}

## 输出格式
返回 JSON 数组，每项包含：
- section: 章节名称（"审核要点"/"问卷要点"/"填报口径微调"/其它原文章节名）
- content: 要点内容（一句话概括）
- common_error: 常见错误（如果有）
- suggest_quiz: 是否建议出题（true/false）

## 规则
1. 只提取需要调查员"记住并应用"的内容
2. 时间安排、通知对象、流程说明等不提取
3. 每个要点聚焦一个知识点
4. 常见错误来自实际填报中的典型误判
5. 目标数量 {keypoint_count} 个：不足时按实际提取数量输出，不凑数、不重复
6. 只输出 JSON，不要 markdown 代码块，不要任何解释"""


def extract_docx_text(path: str) -> str:
    """用 python-docx 读取 docx 全部段落文本（含表格单元格），返回合并文本。"""
    from docx import Document

    doc = Document(path)
    parts: list[str] = []
    for p in doc.paragraphs:
        t = (p.text or "").strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                t = (cell.text or "").strip()
                if t:
                    parts.append(t)
    return "\n".join(parts)


def extract_pdf_text(path: str) -> str:
    """用 pdfplumber 读取 PDF 每页文字层文本（扫描件无文字层 → 按约定不做 OCR）。"""
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = (page.extract_text() or "").strip()
            if t:
                parts.append(t)
    return "\n".join(parts)


def extract_pptx_text(path: str) -> str:
    """用 python-pptx 读取 PPTX 每页文本框 + 表格 + 备注文本（图片内容不提取）。"""
    from pptx import Presentation

    prs = Presentation(path)
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = (shape.text_frame.text or "").strip()
                if t:
                    lines.append(t)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [(c.text or "").strip() for c in row.cells]
                    line = " ".join(x for x in cells if x)
                    if line:
                        lines.append(line)
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                lines.append("【备注】" + notes)
        if lines:
            parts.append(f"第{idx}页：" + "；".join(lines))
    return "\n".join(parts)


def extract_file_text(path: str) -> str:
    """按扩展名提取文件文本：docx→python-docx；doc/wps→COM；pdf→pdfplumber；pptx→python-pptx。"""
    ext = Path(path).suffix.lower()
    if ext == ".docx":
        return extract_docx_text(path)
    if ext == ".pdf":
        return extract_pdf_text(path)
    if ext == ".pptx":
        return extract_pptx_text(path)
    return _extract_legacy_doc_text(path)


def extract_doc_text(path: str) -> str:
    """兼容别名：转发到 extract_file_text（旧引用/测试）。"""
    return extract_file_text(path)


def _try_soffice(path: str) -> str | None:
    """用 LibreOffice(headless) 把 .doc/.wps 转成 txt；不可用/失败返回 None。"""
    import shutil
    import subprocess
    import tempfile

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return None
    try:
        with tempfile.TemporaryDirectory() as outdir:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "txt:Text", "--outdir", outdir, str(path)],
                capture_output=True, timeout=120,
            )
            out_txt = Path(outdir) / (Path(path).stem + ".txt")
            if out_txt.exists():
                text = out_txt.read_text(encoding="utf-8", errors="replace")
                if text.strip():
                    return text.strip().replace("\r", "\n")
    except Exception:
        return None
    return None


def _extract_legacy_doc_text(path: str) -> str:
    """.doc/.wps 优先用 LibreOffice 转文本；Windows 无 soffice 时退回 Word/WPS COM。

    均不可用时抛清晰错误。
    """
    # 1) LibreOffice（跨平台，Linux 服务器推荐）
    text = _try_soffice(path)
    if text:
        return text
    # 2) Word/WPS COM（Windows + Office/WPS）
    try:
        import pythoncom
        from win32com.client import Dispatch
    except ImportError as e:
        raise RuntimeError(
            "无法解析 .doc/.wps：服务器缺少 LibreOffice 且本机无 Word/WPS。请将文件转存为 .docx 再导入。"
        ) from e

    pythoncom.CoInitialize()
    app = None
    try:
        for prog_id in ("Word.Application", "KWPS.Application"):
            try:
                app = Dispatch(prog_id)
                app.Visible = False
                try:
                    app.DisplayAlerts = 0
                except Exception:
                    pass
                doc = app.Documents.Open(str(path), ReadOnly=True)
                try:
                    text = doc.Content.Text
                finally:
                    try:
                        doc.Close(False)
                    except Exception:
                        pass
                return (text or "").replace("\r", "\n").strip()
            except Exception:
                app = None
                continue
        raise RuntimeError(
            "无法解析 .doc/.wps：本机 Word/WPS 均打不开该文件，请转存为 .docx 再导入。"
        )
    finally:
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()


def segment_notice(text: str) -> list[dict]:
    """按章节标题切分通知文本。

    返回 [{section, text}]；标题前的导言归入"其它"。
    """
    segments: list[dict] = []
    current: dict | None = None
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if _is_section_header(line):
            section = KNOWN_SECTIONS.get(_strip_numbering(line), _strip_numbering(line))
            current = {"section": section, "text": line}
            segments.append(current)
        else:
            if current is None:
                current = {"section": "其它", "text": ""}
                segments.append(current)
            current["text"] = (current["text"] + "\n" + line).strip()
    return segments


def strip_json_fence(raw: str) -> str:
    """去掉 markdown 代码块围栏（```json ... ```）。"""
    s = raw.strip()
    s = re.sub(r"^```(?:json)?\s*", "", s, flags=re.IGNORECASE)
    s = re.sub(r"\s*```$", "", s)
    return s.strip()


def parse_llm_json(raw: str):
    """从 LLM 输出中解析 JSON；失败返回 None（由调用方决定重试）。"""
    if not raw:
        return None
    s = strip_json_fence(raw)
    try:
        return json.loads(s)
    except json.JSONDecodeError:
        pass
    # 兜底：截取第一个 [ 或 { 到最后一个 ] 或 } 之间的内容
    start = min([i for i in (s.find("["), s.find("{")) if i != -1], default=-1)
    end = max(s.rfind("]"), s.rfind("}"))
    if start != -1 and end > start:
        try:
            return json.loads(s[start : end + 1])
        except json.JSONDecodeError:
            return None
    return None


def normalize_keypoints(items: list) -> list[dict]:
    """校验并规范化 LLM 要点数组。非法项跳过。"""
    out: list[dict] = []
    for it in items or []:
        if not isinstance(it, dict):
            continue
        content = (it.get("content") or "").strip()
        if not content:
            continue
        section = (it.get("section") or "").strip() or "其它"
        out.append({
            "section": KNOWN_SECTIONS.get(section, section),
            "content": content,
            "common_error": (it.get("common_error") or "").strip(),
            "suggest_quiz": bool(it.get("suggest_quiz", True)),
        })
    return out


def find_source_quote(content: str, section_text: str) -> str:
    """在来源段落里找与要点最相关的一句话作为 source_quote。

    策略：优先返回包含 content 的完整句；否则逐句 difflib 相似度 ≥ 0.35 取最高；
    仍无 → 空串（由前端/管理员兜底）。
    """
    if not section_text:
        return ""
    sentences = [s.strip() for s in re.split(r"[。；\n]", section_text) if s.strip()]
    if content:
        for s in sentences:
            if content in s:
                return s
    best, best_ratio = "", 0.0
    for s in sentences:
        r = SequenceMatcher(None, content, s).ratio()
        if r > best_ratio:
            best, best_ratio = s, r
    return best if best_ratio >= 0.35 else ""


def parse_keypoints(raw: str):
    """LLM 原始输出 → 规范化要点数组；非法 JSON 返回 None（触发重试）。"""
    data = parse_llm_json(raw)
    if not isinstance(data, list):
        return None
    return normalize_keypoints(data)


def _llm_json(messages: list[dict], llm_chat_fn, parse_fn, what: str):
    """带重试的 LLM JSON 调用：失败后追加"只输出 JSON"提示重试。

    网络层异常（超时/限流）自动退避重试，最多 2 次。
    """
    raw = _llm_call_with_retry(llm_chat_fn, messages)
    parsed = parse_fn(raw)
    for _ in range(QUIZ_RETRY_TIMES):
        if parsed is not None:
            return parsed
        repair = messages + [
            {"role": "assistant", "content": raw},
            {"role": "user", "content": "只输出合法 JSON，不要 markdown 代码块，不要任何解释。"},
        ]
        if not raw or not raw.strip():
            repair = messages + [
                {"role": "user", "content": "刚才没有收到输出。请直接输出合法 JSON，不要 markdown 代码块，不要任何解释。"},
            ]
        raw = _llm_call_with_retry(llm_chat_fn, repair)
        parsed = parse_fn(raw)
    if parsed is not None:
        return parsed
    snippet = (raw or "")[:200].replace("\n", " ")
    raise ValueError(f"{what}：LLM 连续返回非法 JSON（最后输出：{snippet}）")




def _llm_call_with_retry(llm_chat_fn, messages: list[dict], retries: int = 2) -> str:
    """调用 LLM；网络层异常（超时/连接/限流）按 2s/4s 退避重试。"""
    import time as _time

    last: Exception | None = None
    for attempt in range(retries + 1):
        try:
            return llm_chat_fn(messages)
        except Exception as e:  # noqa: BLE001
            last = e
            if attempt < retries:
                _time.sleep(2 * (attempt + 1))
    raise last if last else RuntimeError("LLM 调用失败")

def run_extraction(notice_text: str, llm_chat_fn, keypoint_count: int | None = None) -> list[dict]:
    """提取编排：分段 → Prompt1（目标要点数）→ 解析（重试）→ 补 section/source_quote。"""
    target = keypoint_count if keypoint_count is not None else QUIZ_DEFAULT_KEYPOINTS
    target = max(1, min(int(target), QUIZ_MAX_KEYPOINTS))
    segments = segment_notice(notice_text)
    messages = [
        {"role": "system", "content": PROMPT1_SYSTEM},
        {"role": "user", "content": PROMPT1_USER.format(notice_text=notice_text[:12000], keypoint_count=target)},
    ]
    keypoints = _llm_json(messages, llm_chat_fn, parse_keypoints, "要点提取")
    # 为每个要点补 source_quote：优先在其 section 段落里找
    for kp in keypoints:
        seg_text = next(
            (s["text"] for s in segments if s["section"] == kp["section"]), ""
        )
        if not seg_text:
            seg_text = "\n".join(s["text"] for s in segments)
        kp["source_quote"] = find_source_quote(kp["content"], seg_text)
    return keypoints

